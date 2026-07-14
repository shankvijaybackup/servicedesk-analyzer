"""PPTX renderer: turns the slide outline (deliverable M) into a .pptx file
using python-pptx. Returns bytes; nothing written to disk here.
"""

import io

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ACCENT = RGBColor(0x24, 0x56, 0xD6)
INK = RGBColor(0x1A, 0x23, 0x32)
SUB = RGBColor(0x5A, 0x6A, 0x7E)


def render(report: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Title slide
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.4))
    p = tb.text_frame.paragraphs[0]
    p.text = "Service Desk Intelligence Assessment"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    sub = s.shapes.add_textbox(Inches(0.8), Inches(3.9), Inches(11.7), Inches(1.2))
    sp = sub.text_frame.paragraphs[0]
    sp.text = (f"{report['meta']['source_name']}  |  {report['meta']['date_range_str']}  |  "
               f"{report['n_analyzed']} records  |  no data retained")
    sp.font.size = Pt(16)
    sp.font.color.rgb = SUB

    for slide_def in report["slides"][1:]:
        s = prs.slides.add_slide(blank)
        # Title bar
        tb = s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(0.9))
        p = tb.text_frame.paragraphs[0]
        p.text = slide_def["title"]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        # Bullets
        body = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.4))
        tf = body.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(slide_def["bullets"][:8]):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = str(bullet)
            para.font.size = Pt(16)
            para.font.color.rgb = INK
            para.space_after = Pt(10)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
