"""Executive PowerPoint deck (deliverable M), built with python-pptx."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

_INK = RGBColor(0x0F, 0x17, 0x2A)
_MUTED = RGBColor(0x64, 0x74, 0x8B)
_ACCENT = RGBColor(0x4F, 0x46, 0xE5)


def write(a: dict, path) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    meta, dq, ex, opp = a["meta"], a["data_quality"], a["executive"], a["opportunities"]
    roi = opp["roi_summary"]

    # 1. Title
    s = prs.slides.add_slide(blank)
    _text(s, "Service Desk Analysis", Inches(0.7), Inches(2.4), Inches(12), Inches(1.2),
          size=40, bold=True, color=_ACCENT)
    _text(s, f"{meta['source_name']}  |  generated {meta['generated_at']}",
          Inches(0.7), Inches(3.6), Inches(12), Inches(0.6), size=18, color=_MUTED)
    _text(s, "Deterministic analysis. No model, no training, no data retained.",
          Inches(0.7), Inches(4.2), Inches(12), Inches(0.5), size=13, color=_MUTED)

    # 2. Executive summary
    s = _content(prs, blank, "Executive summary")
    _bullets(s, ex["current_state_summary"] + [f"[{f['confidence']}] {f['finding']}"
                                               for f in ex["top_findings"]])

    # 3. Data quality
    s = _content(prs, blank, "Data quality and confidence")
    _bullets(s, [
        f"Records: {dq['total_records']}",
        f"Quality: {dq['quality_grade']} ({dq['quality_score']}/100)",
        f"Missing fields: {', '.join(dq['fields_missing']) or 'none'}",
        f"MTTR: {'available' if dq['mttr_available'] else 'not available'}",
        f"Duplicates: {dq['duplicate_records']}",
    ])

    # 4. Volume and mix
    vol = a["volume"]
    lines = [f"Total tickets: {vol['total']}"]
    if "open_backlog" in vol:
        lines.append(f"Resolved: {vol['resolved']}  |  Open backlog: {vol['open_backlog']}")
    for key, label in [("by_type", "Type"), ("by_priority", "Priority")]:
        if key in vol:
            lines.append(f"{label}: " + ", ".join(f"{r['value']} {r['pct']}%" for r in vol[key][:5]))
    s = _content(prs, blank, "Ticket volume and mix")
    _bullets(s, lines)

    # 5. MTTR
    m = a["mttr"]
    if m.get("available"):
        lines = [f"Overall median {m['overall']['median_hours']}h, p90 {m['overall']['p90_hours']}h"]
        if m.get("slowest"):
            lines.append("Slowest: " + ", ".join(f"{r['value']} {r['median_hours']}h" for r in m["slowest"][:3]))
        if m.get("fastest"):
            lines.append("Fastest: " + ", ".join(f"{r['value']} {r['median_hours']}h" for r in m["fastest"][:3]))
    else:
        lines = [m.get("note", "MTTR not available.")]
    s = _content(prs, blank, "MTTR: slowest and fastest")
    _bullets(s, lines)

    # 6. Themes (Pareto)
    s = _content(prs, blank, "Theme breakdown (the vital few)")
    _bullets(s, [f"{t['theme']}: {t['count']} ({t['pct']}%) [{t['confidence']}]"
                 for t in a["themes"][:8]])

    # 7. Friction points
    s = _content(prs, blank, "Top operational friction points")
    _bullets(s, [f"{fp['area']}: {fp['count']} tickets, {fp['mttr_median_hours']}h median, "
                 f"~{fp['est_effort_hours']}h effort" for fp in ex["friction_points"]])

    # 8. Automation backlog
    s = _content(prs, blank, "AI automation opportunity backlog")
    _bullets(s, [f"{b['theme']}: {b['primary_type']}, {b['est_deflectable_tickets_range'][0]}-"
                 f"{b['est_deflectable_tickets_range'][1]} deflectable [{b['confidence']}]"
                 for b in opp["backlog"][:6]])

    # 9. Agentic
    s = _content(prs, blank, "Agentic AI use cases (with guardrails)")
    if opp["agentic_backlog"]:
        _bullets(s, [f"{g['theme']}: {g['system_of_action']}; approval: {g['human_approval_required']}"
                     for g in opp["agentic_backlog"][:6]])
    else:
        _bullets(s, ["No agentic use cases met the volume threshold in this dataset."])

    # 10. Solution mapping
    s = _content(prs, blank, "Atomicwork solution mapping")
    _bullets(s, [f"{sm['theme']}: {', '.join(sm['capabilities'][:3])}"
                 for sm in opp["solution_map"][:8] if sm["capabilities"]])

    # 11. Roadmap
    s = _content(prs, blank, "30-60-90 day roadmap")
    rm = ex["roadmap_30_60_90"]
    _bullets(s, [
        f"0-30: {rm['days_0_30']['focus']} ({', '.join(rm['days_0_30']['themes']) or 'top themes'})",
        f"30-60: {rm['days_30_60']['focus']} ({', '.join(rm['days_30_60']['themes']) or 'workflow themes'})",
        f"60-90: {rm['days_60_90']['focus']} ({', '.join(rm['days_60_90']['themes']) or 'integration themes'})",
    ])

    # 12. ROI + assumptions
    s = _content(prs, blank, "ROI ranges and assumptions")
    _bullets(s, [
        f"Estimated deflectable: {roi['est_total_deflectable_range'][0]}-{roi['est_total_deflectable_range'][1]} "
        f"tickets ({roi['est_total_deflectable_pct_range'][0]}-{roi['est_total_deflectable_pct_range'][1]}%)",
    ] + roi["assumptions"])

    # 13. Recommendations + questions
    s = _content(prs, blank, "Recommendations and next steps")
    _bullets(s, ex["final_recommendations"] + ["Workshop questions:"] + ex["workshop_questions"][:4])

    # 14. Implementation and UAT package
    impl = a.get("implementation")
    if impl:
        tp = impl["test_plan"]
        s = _content(prs, blank, "Implementation and UAT package")
        lines = [
            f"UAT test plan: {tp['total_cases']} cases "
            f"(Must {tp['must_count']}, Should {tp['should_count']}, Could {tp['could_count']}), "
            f"generated from the observed themes, groups, and priorities",
            f"Role matrix: {len(impl['role_matrix'])} roles including per-group agent access",
            f"Readiness checklist: "
            f"{sum(1 for r in impl['readiness_checklist'] if r['gate'] == 'Must')} Must gates "
            f"before go-live",
            "15-day testing phase plan:",
        ]
        lines += [f"    Days {p['days']}: {p['phase']} (exit: {p['exit_criteria']})"
                  for p in impl["phase_plan_15_day"]]
        _bullets(s, lines, size=14)

    iteration = a.get("iteration")
    if iteration:
        s = _content(prs, blank, "Iterative improvement scorecard")
        lines = [
            f"Pilot: {iteration['pilot']['name']}",
            f"Decision: {iteration['decision']['code'].replace('_', ' ').title()}",
            f"Comparability: {iteration['comparability']['status']}",
        ]
        lines.extend(iteration["decision"]["reasons"])
        lines.append(iteration["causality_note"])
        _bullets(s, lines, size=15)

    path = Path(path)
    prs.save(path)
    return str(path)


def _content(prs, layout, title):
    s = prs.slides.add_slide(layout)
    _text(s, title, Inches(0.6), Inches(0.4), Inches(12), Inches(0.9), size=28, bold=True, color=_ACCENT)
    return s


def _text(slide, text, left, top, width, height, *, size=18, bold=False, color=_INK,
          align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _bullets(slide, items, *, left=Inches(0.7), top=Inches(1.5), width=Inches(12),
             height=Inches(5.4), size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"-  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = _INK
        p.space_after = Pt(6)
