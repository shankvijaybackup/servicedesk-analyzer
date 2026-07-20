from __future__ import annotations

import json

import pandas as pd
from openpyxl import load_workbook
from pptx import Presentation

from sda import CohortSpec, PilotCharter, analyze_dataframe, compare_dataframes
from sda.report import write_all


def _df(prefix: str, hours: int):
    return pd.DataFrame({
        "Number": [f"{prefix}-{i}" for i in range(20)],
        "Opened": ["2026-01-01"] * 20,
        "Resolution time hours": [hours] * 20,
        "Status": ["Resolved"] * 20,
        "Priority": ["P3"] * 20,
        "Assignment group": ["L1"] * 20,
        "Short description": ["PRIVATE-CANARY password reset"] * 20,
    })


def test_all_writers_support_iteration_and_exclude_raw_description(tmp_path):
    baseline, follow = _df("B", 10), _df("F", 5)
    charter = PilotCharter(
        pilot_id="p", name="Access pilot", intervention="Draft",
        cohort=CohortSpec(themes=("Access & Authentication",)),
        minimum_cohort_size=5,
    )
    analysis = analyze_dataframe(follow, source_name="follow-up.csv")
    analysis["iteration"] = compare_dataframes(baseline, follow, charter)
    paths = write_all(analysis, tmp_path)
    assert len(paths) == 5

    assert "PRIVATE-CANARY" not in (tmp_path / "servicedesk-analysis.html").read_text()
    assert "PRIVATE-CANARY" not in (tmp_path / "servicedesk-analysis.md").read_text()
    payload = json.loads((tmp_path / "servicedesk-analysis.json").read_text())
    assert payload["iteration"]["decision"]["code"] == "widen"
    assert "PRIVATE-CANARY" not in json.dumps(payload)

    workbook = load_workbook(tmp_path / "servicedesk-analysis.xlsx", read_only=True)
    assert "IterationScorecard" in workbook.sheetnames
    excel_text = " ".join(str(cell.value) for ws in workbook.worksheets for row in ws for cell in row)
    assert "PRIVATE-CANARY" not in excel_text

    deck = Presentation(tmp_path / "servicedesk-analysis.pptx")
    slide_text = " ".join(shape.text for slide in deck.slides for shape in slide.shapes
                          if hasattr(shape, "text"))
    assert "Iterative improvement scorecard" in slide_text
    assert "PRIVATE-CANARY" not in slide_text
